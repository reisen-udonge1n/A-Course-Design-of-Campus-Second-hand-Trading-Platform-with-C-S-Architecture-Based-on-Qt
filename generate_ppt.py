#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""生成课程设计答辩PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ======================== 颜色方案 ========================
COLOR_BG = RGBColor(0x1A, 0x1A, 0x2E)         # 深蓝黑背景
COLOR_ACCENT = RGBColor(0x00, 0xD2, 0xFF)      # 青色强调
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT = RGBColor(0xCC, 0xCC, 0xDD)
COLOR_CARD_BG = RGBColor(0x25, 0x25, 0x3A)     # 卡片背景
COLOR_GREEN = RGBColor(0x2E, 0xCC, 0x71)
COLOR_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
COLOR_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
COLOR_RED = RGBColor(0xE7, 0x4C, 0x3C)
COLOR_YELLOW = RGBColor(0xF1, 0xC4, 0x0F)

def add_bg(slide, color=COLOR_BG):
    """设置幻灯片背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=""):
    """添加顶部标题栏"""
    # 顶部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_LIGHT

def add_section_number(slide, number, color):
    """添加章节编号圆形"""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.4), Inches(0.2), Inches(0.65), Inches(0.65)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

def add_card(slide, left, top, width, height, title, content, accent_color=COLOR_ACCENT, icon=""):
    """添加卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.fill.background()
    # 圆角
    card.adjustments[0] = 0.05

    # 左侧色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.15), Inches(0.06), height - Inches(0.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # 标题
    header = icon + " " + title if icon else title
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08), width - Inches(0.4), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = accent_color

    # 内容
    txBox2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), height - Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(content.split('\n')):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_LIGHT
        p.space_after = Pt(3)

def add_footer(slide, text="校园二手交易平台 · 课程设计答辩"):
    """添加页脚"""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x88)
    p.alignment = PP_ALIGN.CENTER

# ======================== 封面 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# 装饰元素
for i, (x, color) in enumerate([
    (Inches(1), COLOR_ACCENT),
    (Inches(1.5), COLOR_GREEN),
    (Inches(2), COLOR_ORANGE),
    (Inches(2.5), COLOR_PURPLE),
]):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.5), Inches(1.8), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

# 主标题
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "校园二手交易平台"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE

p2 = tf.add_paragraph()
p2.text = "Campus Second-Hand Trading Platform"
p2.font.size = Pt(22)
p2.font.color.rgb = COLOR_ACCENT
p2.space_before = Pt(8)

# 副标题
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11), Inches(1.0))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "基于 Qt6/C++ 与 MySQL 的桌面端交易平台"
p3.font.size = Pt(20)
p3.font.color.rgb = COLOR_LIGHT

p4 = tf2.add_paragraph()
p4.text = "涵盖用户系统 · 数据库设计 · 社交关注 · 实时聊天 · 界面交互"
p4.font.size = Pt(16)
p4.font.color.rgb = RGBColor(0x88, 0x88, 0xAA)
p4.space_before = Pt(10)

# 底部信息
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11), Inches(0.8))
tf3 = txBox3.text_frame
p5 = tf3.paragraphs[0]
p5.text = "技术栈: Qt 6.10.3 | C++17 | MariaDB/MySQL | TCP/SSL | Qt Widgets"
p5.font.size = Pt(14)
p5.font.color.rgb = RGBColor(0x77, 0x77, 0x99)

add_footer(slide, "")

# ======================== 目录 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "目录", "CONTENTS")

sections = [
    ("01", "注册登录系统", "邮箱验证码注册 · SHA-256登录 · 会话管理 · 密码找回", COLOR_ACCENT),
    ("02", "数据库系统", "MySQL/MariaDB · 6表设计 · Schema Evolution · 参数化查询", COLOR_GREEN),
    ("03", "关注系统", "关注/取消关注 · 用户发现 · 搜索筛选 · 商品查看", COLOR_ORANGE),
    ("04", "聊天系统", "TCP实时通信 · HTTP图片服务 · 在线状态 · 消息通知", COLOR_PURPLE),
    ("05", "页面UI设计", "12页面架构 · 卡片式布局 · 底部导航 · 灵动岛通知", COLOR_RED),
]

for i, (num, title, desc, color) in enumerate(sections):
    y = Inches(1.6 + i * 1.05)
    # 编号
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 标题与描述
    txBox = slide.shapes.add_textbox(Inches(2.4), y - Inches(0.02), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_LIGHT
    p2.space_before = Pt(4)

add_footer(slide)

# ======================== 第一部分：注册登录系统 ========================
# 标题页
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 1, COLOR_ACCENT)
add_title_bar(slide, "第一部分：注册登录系统", "邮箱验证码 · SHA-256认证 · 会话管理 · 密码找回")

# 功能卡片
cards_data = [
    (Inches(0.5), Inches(1.5), Inches(4.0), Inches(2.6),
     "用户注册 RegisterDialog",
     "• 邮箱 + 验证码注册，确保邮箱真实性\n"
     "• 用户名/邮箱唯一性校验，防止重复注册\n"
     "• SHA-256 哈希密码存储，不存明文\n"
     "• 60秒发送冷却计时，防止验证码滥用\n"
     "• 完整表单校验（空值、格式、两次密码一致性）",
     COLOR_ACCENT, "📝"),

    (Inches(4.8), Inches(1.5), Inches(4.0), Inches(2.6),
     "用户登录 LoginWindow",
     "• 支持用户名或邮箱两种方式登录\n"
     "• SHA-256 密码哈希比对验证\n"
     "• 5次登录失败锁定，防暴力破解\n"
     "• 生成64位随机会话令牌，写入数据库\n"
     "• 参数化查询防SQL注入\n"
     "• 回车键快捷登录",
     COLOR_GREEN, "🔑"),

    (Inches(9.1), Inches(1.5), Inches(4.0), Inches(2.6),
     "密码找回 ForgotPasswordDialog",
     "• 两步验证流程：邮箱验证 → 重置密码\n"
     "• 邮箱未注册时引导跳转注册页\n"
     "• 密码长度最低6位校验\n"
     "• 新密码SHA-256哈希后更新数据库",
     COLOR_ORANGE, "🔒"),

    (Inches(0.5), Inches(4.4), Inches(4.0), Inches(2.4),
     "验证码管理 VerificationCodeManager",
     "• Meyer's Singleton 单例模式\n"
     "• QRandomGenerator 生成6位数字验证码\n"
     "• QMap<email, {code, time}> 内存存储\n"
     "• 5分钟过期自动清理（generate时触发）\n"
     "• 验证后立即删除，防重放攻击\n"
     "• 进程重启自动失效（安全特性）",
     COLOR_PURPLE, "📧"),

    (Inches(4.8), Inches(4.4), Inches(4.0), Inches(2.4),
     "SMTP 邮件发送 SmtpSender",
     "• 显式有限状态机驱动 SMTP 会话\n"
     "• Init→EHLO→AUTH LOGIN→USER→PASS→\n"
     "  MAIL FROM→RCPT TO→DATA→BODY→QUIT\n"
     "• SSL/TLS 加密连接（端口465）\n"
     "• Base64 编码认证信息\n"
     "• 10秒超时保护，防止连接挂死\n"
     "• QSslSocket 非阻塞异步发送",
     COLOR_RED, "✉️"),

    (Inches(9.1), Inches(4.4), Inches(4.0), Inches(2.4),
     "会话安全机制",
     "• 登录时生成 session_token 写入 users 表\n"
     "• 新登录覆盖旧令牌（单设备登录策略）\n"
     "• MainWindow 每30秒校验 session_token\n"
     "• 其他设备登录时自动退出并提示\n"
     "• 主动退出清空 token，防止误判",
     COLOR_YELLOW, "🛡️"),
]
for args in cards_data:
    add_card(slide, *args)

add_footer(slide)

# ======================== 第二部分：数据库系统 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 2, COLOR_GREEN)
add_title_bar(slide, "第二部分：数据库系统", "MySQL/MariaDB · 6表设计 · Schema Evolution · 参数化查询")

db_cards = [
    (Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.8),
     "DbManager — 数据库自举与 Schema Evolution",
     "• Meyer's Singleton 单例模式，进程级唯一连接\n"
     "• 双阶段初始化：\n"
     "  ① 无库临时连接 → CREATE DATABASE IF NOT EXISTS\n"
     "  ② 目标库正式连接 → CREATE TABLE IF NOT EXISTS\n"
     "• Schema Evolution 策略：ALTER TABLE ADD COLUMN\n"
     "  + 忽略 1060 错误码实现幂等性\n"
     "• SET FOREIGN_KEY_CHECKS 管理外键依赖\n"
     "• utf8mb4 字符集，支持完整 Unicode（含 Emoji）\n"
     "• 所有查询使用参数化绑定，防止 SQL 注入",
     COLOR_GREEN, "🗄️"),

    (Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8),
     "6张核心表设计",
     "• users — 用户账户（id, username, email, password_hash, session_token）\n"
     "• products — 商品信息（id, title, price, description, seller_id, status,\n"
     "  category, images LONGTEXT, created_at）+ 外键→users\n"
     "• orders — 订单信息（product_id, buyer_id, seller_id, status,\n"
     "  offer_price, created_at）+ 三外键约束\n"
     "• messages — 聊天消息（sender_id, receiver_id, product_id,\n"
     "  content, is_read, created_at）+ 两外键→users\n"
     "• user_profiles — 用户资料（user_id PK+FK, gender, address,\n"
     "  shipping/receiving_address, avatar TEXT）\n"
     "• followers — 关注关系（follower_id, followed_id,\n"
     "  created_at, UNIQUE联合唯一约束）+ 两外键→users",
     COLOR_ACCENT, "📊"),
]

for args in db_cards:
    add_card(slide, *args)

# 下半部分
bottom_cards = [
    (Inches(0.5), Inches(4.6), Inches(6.0), Inches(2.5),
     "图片存储方案",
     "• 两种存储路径并存：\n"
     " ① Base64 Data URI — 直接存入 products.images LONGTEXT\n"
     " ② HTTP URL — 上传到 ChatServer 的 images/ 目录\n"
     "• loadPixmap() 统一加载：自动识别 data: / http:// / 本地路径\n"
     "• QNetworkAccessManager + QEventLoop 同步HTTP加载\n"
     "• 内存缓存 static QMap，避免重复请求\n"
     "• 15秒超时保护",
     COLOR_ORANGE, "🖼️"),

    (Inches(6.8), Inches(4.6), Inches(6.0), Inches(2.5),
     "订单生命周期管理",
     "• 6种订单状态流转：\n"
     "  0=PendingConfirm → 2=PendingShip → 5=Shipped → 3=Completed\n"
     "  任意状态可转为 4=Cancelled\n"
     "• offer_price 支持议价功能（nullable DECIMAL）\n"
     "• 买家 → 卖家 → 买家 双向状态流转\n"
     "• ProfilePage 6维统计：发布/购买/待确认/待发货/待收货/已完成\n"
     "• 6个 OrderListPage 实例，按不同 SQL 条件查询",
     COLOR_PURPLE, "📋"),
]
for args in bottom_cards:
    add_card(slide, *args)

add_footer(slide)

# ======================== 第三部分：关注系统 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 3, COLOR_ORANGE)
add_title_bar(slide, "第三部分：关注系统", "关注/取消关注 · 用户发现 · 搜索筛选 · 商品查看")

cards_follow = [
    (Inches(0.5), Inches(1.5), Inches(4.0), Inches(2.6),
     "FollowersDialog — 双标签页设计",
     "• Tab 1「已关注」— JOIN查询当前用户的关注列表\n"
     "  按 created_at 倒序，支持用户名搜索\n"
     "  「取消关注」按钮，DELETE 操作 + 确认弹窗\n"
     "• Tab 2「发现用户」— 列出所有非己用户\n"
     "  已关注用户显示 ✅ 标记且不可重复选择\n"
     "  搜索过滤 + 一键关注 INSERT 操作\n"
     "• 按钮样式动态切换（蓝/红）",
     COLOR_ORANGE, "👥"),

    (Inches(4.8), Inches(1.5), Inches(4.0), Inches(2.6),
     "数据库层 followers 表设计",
     "• UNIQUE(follower_id, followed_id) 联合唯一约束\n"
     "  防止重复关注\n"
     "• follower_id / followed_id 双外键 → users(id)\n"
     "• created_at 记录关注时间\n"
     "• 参数化 SQL 防注入\n"
     "  INSERT/DELETE + SELECT COUNT 状态查询",
     COLOR_GREEN, "🗃️"),

    (Inches(9.1), Inches(1.5), Inches(4.0), Inches(2.6),
     "集成与联动",
     "• HomePage 统计栏显示关注数量\n"
     "• ChatDialog 内嵌关注/取消关注按钮\n"
     "  isFollowing() 实时查询状态 → updateFollowButton()\n"
     "• UserProductsDialog 查看关注用户的商品\n"
     "• 关注变更后 emit followChanged() 刷新统计\n"
     "• 双击列表项查看该用户发布的商品",
     COLOR_PURPLE, "🔗"),

    (Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.5),
     "关注系统交互流程",
     "① 用户点击「👥 关注」→ FollowersDialog 弹窗 → 默认显示「已关注」Tab → loadFollowers() 查询 followers JOIN users\n"
     "② 切换到「发现用户」Tab → loadUsers() 查询所有非己用户 → 已关注的标记 ✅ + 不可选中\n"
     "③ 选中用户 → 点击「关注」→ INSERT INTO followers → 更新UI为 ✅ (已关注) + 不可选 → emit followChanged()\n"
     "④ 选中已关注用户 → 点击「取消关注」→ 确认弹窗 → DELETE FROM followers → 更新UI → emit followChanged()\n"
     "⑤ 双击用户 → emit viewUserProducts() → 打开 UserProductsDialog 查看该用户的在售商品",
     COLOR_ACCENT, "🔄"),
]
for args in cards_follow:
    add_card(slide, *args)

add_footer(slide)

# ======================== 第四部分：聊天系统 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 4, COLOR_PURPLE)
add_title_bar(slide, "第四部分：聊天系统", "TCP实时通信 · HTTP图片服务 · 自定义协议 · 消息通知")

chat_cards_page1 = [
    (Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.6),
     "ChatServer — 双协议服务端 (chat_server.exe)",
     "TCP聊天协议（端口8888）：\n"
     "• 4字节大端序长度前缀 + JSON payload 组帧\n"
     "• 消息类型: login / logout / message / ping / pong / user_list\n"
     "• 单线程事件驱动（Qt信号槽 + 非阻塞I/O）\n"
     "• m_clients 管理所有连接，m_userToSocket 路由消息\n"
     "• 多点登录互踢：新登录→旧连接发送提示→断开\n"
     "• broadcastUserList() 广播在线用户变更\n"
     "• 30s ping/pong 心跳保活",
     COLOR_PURPLE, "🖥️"),

    (Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.6),
     "ChatServer — HTTP 图片服务（端口9999）",
     "• POST /upload — 接收图片数据 → 保存到 images/ 目录\n"
     "  生成唯一文件名: 时间戳_随机数.ext\n"
     "  返回 JSON {url: \"http://host:9999/images/xxx.jpg\"}\n"
     "• GET /images/<file> — 返回图片文件\n"
     "  MIME类型自动识别（jpg/png/gif/bmp/webp）\n"
     "  安全校验：禁止 .. 路径穿越\n"
     "• 手动解析 HTTP 请求头 + Content-Length\n"
     "• 短连接模式（Connection: close）",
     COLOR_GREEN, "🌐"),

    (Inches(0.5), Inches(4.4), Inches(4.0), Inches(2.5),
     "ChatClient — TCP客户端 Singleton",
     "• 持久TCP连接 + JSON协议通信\n"
     "• 自动重连：非主动断开时 5s 间隔重连\n"
     "• 30s ping 心跳，保活检测\n"
     "• 帧同步：4字节长度前缀粘包/拆包处理\n"
     "• QSet<int> 本地维护在线用户状态\n"
     "• 信号驱动：messageReceived/userListChanged\n"
     "  loginSuccess/disconnected",
     COLOR_ORANGE, "📡"),

    (Inches(4.8), Inches(4.4), Inches(4.0), Inches(2.5),
     "ChatBubble — 聊天气泡控件",
     "• 独立 QWidget，轻量渲染\n"
     "• isSelf=true → 蓝色气泡右对齐\n"
     "• isSelf=false → 白色气泡左对齐\n"
     "• 最大宽度260px，自动换行\n"
     "• 时间戳显示（HH:mm格式）\n"
     "• QHBoxLayout + QSpacerItem 弹性对齐\n"
     "• QSS 内联样式（圆角、阴影、配色）",
     COLOR_ACCENT, "💬"),

    (Inches(9.1), Inches(4.4), Inches(4.0), Inches(2.5),
     "消息通知系统（双层设计）",
     "① MainWindow 灵动岛通知栏：\n"
     "  QFrame 顶部胶囊形通知，5s自动隐藏\n"
     "  有打开聊天时不弹出（智能去重）\n"
     "  点击通知直达聊天页\n"
     "② NotificationPopup 浮动弹窗：\n"
     "  居中显示，5s自动隐藏\n"
     "  mousePressEvent 点击跳转\n"
     "  作为无聊天窗口时的兜底通知",
     COLOR_YELLOW, "🔔"),
]
for args in chat_cards_page1:
    add_card(slide, *args)

add_footer(slide)

# 聊天系统第二页
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 4, COLOR_PURPLE)
add_title_bar(slide, "第四部分：聊天系统（续）", "ChatDialog · ChatListPage · 消息持久化 · 敏感词过滤")

chat_cards_page2 = [
    (Inches(0.5), Inches(1.5), Inches(4.0), Inches(2.6),
     "ChatDialog — 模态聊天窗口",
     "• 显示对方在线状态（绿/灰实时更新）\n"
     "• loadMessageHistory() 从DB加载历史消息\n"
     "• 双向查询：(sender=A AND receiver=B) OR 反向\n"
     "• sendMessage() 双路径持久化：\n"
     "  ① INSERT INTO messages（数据库）\n"
     "  ② ChatClient::sendMessage（TCP实时）\n"
     "• 发送前 Filter 敏感词检查\n"
     "• 内嵌关注/取消关注按钮",
     COLOR_PURPLE, "💬"),

    (Inches(4.8), Inches(1.5), Inches(4.0), Inches(2.6),
     "ChatListPage — 会话列表页",
     "• 子查询聚合最近联系人 + 最后消息\n"
     "• 卡片式列表：头像首字母 + 用户名 + 时间 + 预览\n"
     "• 已读/未读样式区分（灰色/深色）\n"
     "• QStackedWidget 内联聊天：\n"
     "  列表页(index 0) ↔ 聊天页(index 1)\n"
     "• 点击进入 → 标记消息已读 → 加载消息\n"
     "• 实时接收：onMessageReceived() 动态插入气泡",
     COLOR_GREEN, "📋"),

    (Inches(9.1), Inches(1.5), Inches(4.0), Inches(2.6),
     "Filter — 敏感词过滤 Singleton",
     "• 硬编码约100条中文敏感词\n"
     "  色情类 / 辱骂类 / 违法类 / 交易类\n"
     "• 子串匹配 contains() 暴力搜索\n"
     "• 大小写不敏感（toLower）\n"
     "• 用于发布商品标题/描述校验\n"
     "  以及发送聊天消息前校验\n"
     "• 改进方向：Aho-Corasick自动机",
     COLOR_RED, "🚫"),

    (Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.5),
     "聊天协议帧格式",
     "┌──────────────────────────────────────┐\n"
     "│  4 bytes (Big-Endian uint32)         │  ← 负载长度 L\n"
     "├──────────────────────────────────────┤\n"
     "│  JSON payload (L bytes, UTF-8)       │  ← QJsonDocument::Compact\n"
     "└──────────────────────────────────────┘\n\n"
     "JSON 消息类型示例：\n"
     "  {\"type\":\"login\", \"user_id\":1, \"username\":\"张三\"}\n"
     "  {\"type\":\"message\", \"to\":2, \"content\":\"你好\"}\n"
     "  {\"type\":\"ping\"} → {\"type\":\"pong\"}\n"
     "  {\"type\":\"user_list\", \"users\":[...]}",
     COLOR_ACCENT, "📦"),

    (Inches(6.8), Inches(4.4), Inches(6.0), Inches(2.5),
     "聊天双路径持久化策略",
     "发送消息时的两条并行路径：\n\n"
     "路径① 数据库持久化：\n"
     "  INSERT INTO messages (sender_id, receiver_id, content, created_at)\n"
     "  → 保证消息不丢失，离线可查看\n\n"
     "路径② TCP实时推送：\n"
     "  ChatClient::sendMessage() → 4字节帧 → ChatServer\n"
     "  → processMessage() → 查找目标socket → sendJson() 转发\n"
     "  → 对方 ChatClient::onReadyRead() → messageReceived 信号\n"
     "  → ChatDialog/ChatListPage 实时显示气泡",
     COLOR_ORANGE, "🔄"),
]
for args in chat_cards_page2:
    add_card(slide, *args)

add_footer(slide)

# ======================== 第五部分：页面UI ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 5, COLOR_RED)
add_title_bar(slide, "第五部分：页面UI设计", "12页面架构 · 卡片式布局 · 底部导航 · 灵动岛通知")

ui_cards = [
    (Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.6),
     "MainWindow — 主框架架构",
     "• QStackedWidget 承载12个页面（0-11索引）\n"
     "• 5按钮底部导航栏：🏠首页 📦商品 ➕发布 🔍搜索 👤我的\n"
     "• QButtonGroup + 互斥选中 + 样式切换\n"
     "• switchPage() 统一页面切换入口，自动刷新数据\n"
     "• 固定窗口尺寸 400×750（移动端比例）\n"
     "• Fusion 风格 + 全局 QSS 样式\n"
     "• 灵动岛通知栏：QFrame 胶囊形，5s自动隐藏",
     COLOR_RED, "🏗️"),

    (Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.6),
     "HomePage — 首页设计",
     "• 欢迎头部：👋 欢迎回来，{用户名}\n"
     "• 三栏统计条：已发布 📦 / 已购买 🛒 / 关注 👥\n"
     "  数据来自 products/orders/followers 表 COUNT\n"
     "  超过99显示为 \"99+\"\n"
     "• 13个分类标签：全部/数码/穿搭/家电/图书/\n"
     "  美妆/运动/食品/宠物/家具/文具/乐器/其他\n"
     "  选中状态红色高亮，点击切换筛选\n"
     "• 热门推荐 Grid：2列卡片布局，每页8个\n"
     "  滚动懒加载，offset分页，防重复加载",
     COLOR_GREEN, "🏠"),

    (Inches(0.5), Inches(4.4), Inches(4.0), Inches(2.5),
     "ProfilePage — 个人中心",
     "• 圆形头像：用户名首字母 + 蓝色背景\n"
     "• 6维订单统计：发布/购买/待确认/待发货/待收货/完成\n"
     "  点击数字跳转对应OrderListPage\n"
     "• 功能菜单：编辑个人资料/我的消息/\n"
     "  我的发布/我的订单/退出登录\n"
     "• 退出→清空session_token→断开聊天→回登录页",
     COLOR_ORANGE, "👤"),

    (Inches(4.8), Inches(4.4), Inches(4.0), Inches(2.5),
     "商品相关页面",
     "• ProductsPage — 双列商品卡片Grid\n"
     "• PublishPage — 发布表单 + 我的发布列表\n"
     "  Filter敏感词校验 → 插入 products 表\n"
     "  publishSuccess 信号自动刷新首页/商品页\n"
     "• SearchPage — LIKE 模糊搜索 + 去重\n"
     "  GROUP BY title 避免重复商品\n"
     "• ProductDetailDialog — 详情 + 议价购买 + 聊天",
     COLOR_ACCENT, "📦"),

    (Inches(9.1), Inches(4.4), Inches(4.0), Inches(2.5),
     "卡片式商品组件",
     "• 统一的 createProductCard() 函数\n"
     "  175×220 固定尺寸，白色圆角卡片\n"
     "• 图片区域（175×130）+ 分类标签 + 标题 + 价格\n"
     "• eventFilter() 处理点击事件\n"
     "  property(\"productId\") 存储商品ID\n"
     "• Hover 效果：背景色微变\n"
     "• 在 HomePage / SearchPage / ProductsPage\n"
     "  三处复用同一组件设计",
     COLOR_PURPLE, "🎴"),
]
for args in ui_cards:
    add_card(slide, *args)

add_footer(slide)

# ======================== 技术总结 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "技术总结", "TECHNICAL SUMMARY")

summary_cards = [
    (Inches(0.5), Inches(1.5), Inches(4.0), Inches(2.6),
     "架构特点",
     "• 客户端-服务器分离架构\n"
     "• 两个独立进程：GUI客户端 + ChatServer\n"
     "• Singleton 模式管理核心服务\n"
     "  DbManager / ChatClient / Filter / VerificationCodeManager\n"
     "• QStackedWidget 单页应用式导航\n"
     "• 信号槽驱动的事件通信\n"
     "• 无ORM，纯参数化SQL",
     COLOR_ACCENT, "🏛️"),

    (Inches(4.8), Inches(1.5), Inches(4.0), Inches(2.6),
     "安全机制",
     "• SHA-256 密码哈希存储\n"
     "• 参数化查询防 SQL 注入\n"
     "• 5次登录失败锁定 + 会话令牌管理\n"
     "• 验证码 5分钟过期 + 一次有效\n"
     "• 敏感词内容过滤（4类约100词）\n"
     "• HTTP 路径遍历防护（.. 检测）\n"
     "• 单设备登录互踢策略",
     COLOR_RED, "🔐"),

    (Inches(9.1), Inches(1.5), Inches(4.0), Inches(2.6),
     "技术栈",
     "• 语言：C++17\n"
     "• 框架：Qt 6.10.3 (Widgets/Sql/Network)\n"
     "• 数据库：MariaDB/MySQL + QMARIADB\n"
     "• 通信：TCP + 自定义JSON协议\n"
     "• 邮件：SMTP/SSL (126.com:465)\n"
     "• 图片：Base64 Data URI + HTTP文件服务\n"
     "• 构建：CMake + MinGW",
     COLOR_GREEN, "⚙️"),

    (Inches(0.5), Inches(4.4), Inches(4.0), Inches(2.5),
     "代码规模",
     "• 约 20+ 个 .h/.cpp 源文件\n"
     "• 6 张数据库表，30+ 字段\n"
     "• 12 个主页面 + 多个对话框\n"
     "• TCP 协议 7 种消息类型\n"
     "• HTTP 2 个路由端点\n"
     "• SMTP 10 状态有限状态机\n"
     "• 约 100 条敏感词过滤库",
     COLOR_PURPLE, "📊"),

    (Inches(4.8), Inches(4.4), Inches(4.0), Inches(2.5),
     "亮点功能",
     "• 实时TCP聊天 + 离线消息DB双路径持久化\n"
     "• 灵动岛式通知栏（类似iOS设计）\n"
     "• 自定义长度前缀JSON通信协议\n"
     "• 数据库零配置自举（auto create DB + tables）\n"
     "• Schema Evolution 向后兼容升级\n"
     "• 订单6状态完整生命周期管理\n"
     "• 价格议价机制（offer_price）",
     COLOR_YELLOW, "⭐"),

    (Inches(9.1), Inches(4.4), Inches(4.0), Inches(2.5),
     "改进方向",
     "• 密码加盐（bcrypt/Argon2）\n"
     "• 敏感词过滤改用 Aho-Corasick 自动机\n"
     "• 验证码迁移至 Redis 存储\n"
     "• 聊天支持图片/文件传输\n"
     "• 指数退避重连策略\n"
     "• HTTPS + WSS 加密通信\n"
     "• 单元测试 + CI/CD 集成",
     COLOR_ORANGE, "🔮"),
]
for args in summary_cards:
    add_card(slide, *args)

add_footer(slide)

# ======================== 结束页 ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 装饰
for i, (x, color) in enumerate([
    (Inches(4.0), COLOR_ACCENT),
    (Inches(4.5), COLOR_GREEN),
    (Inches(5.0), COLOR_ORANGE),
    (Inches(5.5), COLOR_PURPLE),
    (Inches(6.0), COLOR_RED),
]):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.0), Inches(1.0), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "感谢聆听"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "THANK YOU"
p2.font.size = Pt(24)
p2.font.color.rgb = COLOR_ACCENT
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(0.5))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "校园二手交易平台 · 基于 Qt6/C++ 与 MySQL"
p3.font.size = Pt(16)
p3.font.color.rgb = COLOR_LIGHT
p3.alignment = PP_ALIGN.CENTER

# 保存
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "课程设计答辩PPT.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
